from pptx import Presentation
from pptx.util import Inches, Pt

def create_business_case_presentation():
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Business Value Analysis"
    subtitle.text = "Transforming Operations Through Automation"
    
    # Slide 2: Key Performance Metrics
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Key Performance Metrics"
    tf = body.text_frame
    tf.text = "Performance Highlights:"
    
    p = tf.add_paragraph()
    p.text = "- 89% Reduction in Processing Time"
    p = tf.add_paragraph()
    p.text = "- 97.5% Accuracy Rate"
    p = tf.add_paragraph()
    p.text = "- 55% Cost Reduction"
    p = tf.add_paragraph()
    p.text = "- 450% Increased Process Throughput"
    
    # Slide 3: ROI Analysis
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    
    title.text = "Return on Investment (ROI)"
    tf = body.text_frame
    tf.text = "Financial Benefits:"
    
    p = tf.add_paragraph()
    p.text = "- Annual Cost Savings: $350,000"
    p = tf.add_paragraph()
    p.text = "- Implementation Investment: $250,000"
    p = tf.add_paragraph()
    p.text = "- ROI (3 Years): 140%"
    p = tf.add_paragraph()
    p.text = "- Payback Period: 0.7 Years"
    
    # Save Presentation
    prs.save('/Users/rahulmehta/Desktop/Test/demo/metrics/reports/business_case_presentation.pptx')

if __name__ == "__main__":
    create_business_case_presentation()