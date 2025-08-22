"""
Multi-Format Export Manager

Provides comprehensive data export capabilities supporting multiple formats,
compression, encryption, and cloud storage integration.
"""

import asyncio
import json
import xml.etree.ElementTree as ET
import csv
import io
import gzip
import zipfile
import tarfile
from datetime import datetime
from typing import Dict, List, Any, Optional, Union, BinaryIO
from dataclasses import dataclass
from enum import Enum
import pandas as pd
import numpy as np
import openpyxl
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from openpyxl.drawing.image import Image as ExcelImage
import xlsxwriter
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import pptx
from pptx.util import Inches
import boto3
from google.cloud import storage as gcs
import paramiko
import ftplib
import requests
from cryptography.fernet import Fernet
import base64


class ExportFormat(Enum):
    """Supported export formats"""
    CSV = "csv"
    EXCEL = "excel"
    JSON = "json"
    XML = "xml"
    PDF = "pdf"
    POWERPOINT = "powerpoint"
    PARQUET = "parquet"
    AVRO = "avro"
    ORC = "orc"
    HTML = "html"
    MARKDOWN = "markdown"
    TXT = "txt"


class CompressionType(Enum):
    """Compression types"""
    NONE = "none"
    GZIP = "gzip"
    ZIP = "zip"
    TAR = "tar"
    BZIP2 = "bzip2"


class StorageDestination(Enum):
    """Storage destinations"""
    LOCAL = "local"
    S3 = "s3"
    GCS = "gcs"
    AZURE = "azure"
    FTP = "ftp"
    SFTP = "sftp"
    HTTP = "http"


@dataclass
class ExportConfig:
    """Export configuration"""
    format: ExportFormat
    filename: str
    compression: CompressionType = CompressionType.NONE
    encryption: bool = False
    destination: StorageDestination = StorageDestination.LOCAL
    destination_config: Optional[Dict[str, Any]] = None
    metadata: Optional[Dict[str, Any]] = None
    quality_settings: Optional[Dict[str, Any]] = None


@dataclass
class ExportResult:
    """Export operation result"""
    success: bool
    file_path: str
    file_size: int
    format: ExportFormat
    compression: CompressionType
    encrypted: bool
    export_time: datetime
    duration_seconds: float
    row_count: Optional[int] = None
    error_message: Optional[str] = None
    checksum: Optional[str] = None


class ExportManager:
    """
    Advanced export manager with multi-format support
    """
    
    def __init__(self, encryption_key: Optional[str] = None):
        self.encryption_key = encryption_key
        if encryption_key:
            self.cipher = Fernet(encryption_key.encode() if isinstance(encryption_key, str) else encryption_key)
        else:
            self.cipher = None
        
        # Cloud storage clients
        self.s3_client = None
        self.gcs_client = None
        self.azure_client = None
        
        # Default styling for formatted exports
        self.default_styles = {
            'excel': {
                'header_font': Font(bold=True, color="FFFFFF"),
                'header_fill': PatternFill(start_color="366092", end_color="366092", fill_type="solid"),
                'border': Border(left=Side(style='thin'), right=Side(style='thin'), 
                               top=Side(style='thin'), bottom=Side(style='thin'))
            },
            'pdf': {
                'title_style': ParagraphStyle('Title', fontSize=16, spaceAfter=12, alignment=1),
                'header_style': ParagraphStyle('Header', fontSize=12, spaceAfter=6, textColor=colors.darkblue),
                'body_style': ParagraphStyle('Body', fontSize=10, spaceAfter=4)
            }
        }
    
    async def export_data(self, 
                         data: Union[pd.DataFrame, List[Dict], Dict],
                         config: ExportConfig) -> ExportResult:
        """
        Export data in specified format
        
        Args:
            data: Data to export (DataFrame, list of dicts, or dict)
            config: Export configuration
        
        Returns:
            Export result with file information
        """
        start_time = datetime.utcnow()
        
        try:
            # Convert data to DataFrame if needed
            df = self._prepare_dataframe(data)
            
            # Generate file path
            file_path = self._generate_file_path(config)
            
            # Export based on format
            if config.format == ExportFormat.CSV:
                await self._export_csv(df, file_path, config)
            elif config.format == ExportFormat.EXCEL:
                await self._export_excel(df, file_path, config)
            elif config.format == ExportFormat.JSON:
                await self._export_json(data, file_path, config)
            elif config.format == ExportFormat.XML:
                await self._export_xml(data, file_path, config)
            elif config.format == ExportFormat.PDF:
                await self._export_pdf(df, file_path, config)
            elif config.format == ExportFormat.POWERPOINT:
                await self._export_powerpoint(df, file_path, config)
            elif config.format == ExportFormat.PARQUET:
                await self._export_parquet(df, file_path, config)
            elif config.format == ExportFormat.HTML:
                await self._export_html(df, file_path, config)
            elif config.format == ExportFormat.MARKDOWN:
                await self._export_markdown(df, file_path, config)
            else:
                raise ValueError(f"Unsupported export format: {config.format}")
            
            # Apply compression if specified
            if config.compression != CompressionType.NONE:
                file_path = await self._compress_file(file_path, config.compression)
            
            # Apply encryption if specified
            if config.encryption and self.cipher:
                file_path = await self._encrypt_file(file_path)
            
            # Upload to destination if not local
            if config.destination != StorageDestination.LOCAL:
                file_path = await self._upload_to_destination(file_path, config)
            
            # Calculate file size and checksum
            file_size = await self._get_file_size(file_path)
            checksum = await self._calculate_checksum(file_path)
            
            duration = (datetime.utcnow() - start_time).total_seconds()
            
            return ExportResult(
                success=True,
                file_path=file_path,
                file_size=file_size,
                format=config.format,
                compression=config.compression,
                encrypted=config.encryption,
                export_time=start_time,
                duration_seconds=duration,
                row_count=len(df) if isinstance(data, pd.DataFrame) else None,
                checksum=checksum
            )
            
        except Exception as e:
            duration = (datetime.utcnow() - start_time).total_seconds()
            
            return ExportResult(
                success=False,
                file_path="",
                file_size=0,
                format=config.format,
                compression=config.compression,
                encrypted=config.encryption,
                export_time=start_time,
                duration_seconds=duration,
                error_message=str(e)
            )
    
    async def bulk_export(self, 
                         datasets: Dict[str, Union[pd.DataFrame, List[Dict], Dict]],
                         configs: Dict[str, ExportConfig]) -> Dict[str, ExportResult]:
        """
        Export multiple datasets in parallel
        
        Args:
            datasets: Dictionary of dataset names to data
            configs: Dictionary of dataset names to export configs
        
        Returns:
            Dictionary of dataset names to export results
        """
        tasks = []
        for dataset_name, data in datasets.items():
            if dataset_name in configs:
                task = self.export_data(data, configs[dataset_name])
                tasks.append((dataset_name, task))
        
        results = {}
        for dataset_name, task in tasks:
            result = await task
            results[dataset_name] = result
        
        return results
    
    async def _export_csv(self, df: pd.DataFrame, file_path: str, config: ExportConfig):
        """Export DataFrame as CSV"""
        quality_settings = config.quality_settings or {}
        
        df.to_csv(
            file_path,
            index=quality_settings.get('include_index', False),
            encoding=quality_settings.get('encoding', 'utf-8'),
            sep=quality_settings.get('separator', ','),
            quoting=quality_settings.get('quoting', csv.QUOTE_MINIMAL),
            date_format=quality_settings.get('date_format', '%Y-%m-%d'),
            float_format=quality_settings.get('float_format', '%.6f')
        )
    
    async def _export_excel(self, df: pd.DataFrame, file_path: str, config: ExportConfig):
        """Export DataFrame as Excel with advanced formatting"""
        quality_settings = config.quality_settings or {}
        
        with pd.ExcelWriter(file_path, engine='openpyxl') as writer:
            # Write main data
            sheet_name = quality_settings.get('sheet_name', 'Data')
            df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            # Apply formatting
            workbook = writer.book
            worksheet = writer.sheets[sheet_name]
            
            # Header formatting
            if quality_settings.get('format_headers', True):
                for cell in worksheet[1]:
                    cell.font = self.default_styles['excel']['header_font']
                    cell.fill = self.default_styles['excel']['header_fill']
                    cell.alignment = Alignment(horizontal='center')
            
            # Add borders
            if quality_settings.get('add_borders', True):
                for row in worksheet.iter_rows():
                    for cell in row:
                        cell.border = self.default_styles['excel']['border']
            
            # Auto-adjust column widths
            if quality_settings.get('auto_width', True):
                for column in worksheet.columns:
                    max_length = 0
                    column_letter = column[0].column_letter
                    for cell in column:
                        try:
                            if len(str(cell.value)) > max_length:
                                max_length = len(str(cell.value))
                        except:
                            pass
                    adjusted_width = min(max_length + 2, 50)
                    worksheet.column_dimensions[column_letter].width = adjusted_width
            
            # Add charts if specified
            if quality_settings.get('add_charts', False):
                await self._add_excel_charts(worksheet, df, quality_settings)
    
    async def _export_json(self, data: Any, file_path: str, config: ExportConfig):
        """Export data as JSON"""
        quality_settings = config.quality_settings or {}
        
        # Convert DataFrame to dict if needed
        if isinstance(data, pd.DataFrame):
            orient = quality_settings.get('orient', 'records')
            json_data = data.to_dict(orient=orient)
        else:
            json_data = data
        
        with open(file_path, 'w', encoding='utf-8') as f:
            json.dump(
                json_data,
                f,
                indent=quality_settings.get('indent', 2),
                ensure_ascii=quality_settings.get('ensure_ascii', False),
                default=str  # Handle datetime and other non-serializable types
            )
    
    async def _export_xml(self, data: Any, file_path: str, config: ExportConfig):
        """Export data as XML"""
        quality_settings = config.quality_settings or {}
        root_name = quality_settings.get('root_element', 'data')
        item_name = quality_settings.get('item_element', 'item')
        
        root = ET.Element(root_name)
        
        if isinstance(data, pd.DataFrame):
            for _, row in data.iterrows():
                item = ET.SubElement(root, item_name)
                for col, value in row.items():
                    element = ET.SubElement(item, str(col))
                    element.text = str(value) if value is not None else ""
        elif isinstance(data, list):
            for item_data in data:
                item = ET.SubElement(root, item_name)
                if isinstance(item_data, dict):
                    for key, value in item_data.items():
                        element = ET.SubElement(item, str(key))
                        element.text = str(value) if value is not None else ""
        
        tree = ET.ElementTree(root)
        tree.write(file_path, encoding='utf-8', xml_declaration=True)
    
    async def _export_pdf(self, df: pd.DataFrame, file_path: str, config: ExportConfig):
        """Export DataFrame as PDF with formatting"""
        quality_settings = config.quality_settings or {}
        
        doc = SimpleDocTemplate(file_path, pagesize=A4)
        elements = []
        styles = getSampleStyleSheet()
        
        # Add title
        title = quality_settings.get('title', 'Data Export')
        title_para = Paragraph(title, self.default_styles['pdf']['title_style'])
        elements.append(title_para)
        elements.append(Spacer(1, 12))
        
        # Add metadata
        if config.metadata:
            for key, value in config.metadata.items():
                meta_para = Paragraph(f"<b>{key}:</b> {value}", styles['Normal'])
                elements.append(meta_para)
            elements.append(Spacer(1, 12))
        
        # Convert DataFrame to table
        table_data = [df.columns.tolist()] + df.values.tolist()
        
        # Create table with styling
        table = Table(table_data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black)
        ]))
        
        elements.append(table)
        doc.build(elements)
    
    async def _export_powerpoint(self, df: pd.DataFrame, file_path: str, config: ExportConfig):
        """Export DataFrame as PowerPoint presentation"""
        quality_settings = config.quality_settings or {}
        
        prs = pptx.Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = quality_settings.get('title', 'Data Export')
        subtitle.text = quality_settings.get('subtitle', f'Generated on {datetime.now().strftime("%Y-%m-%d")}')
        
        # Data slide
        bullet_slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = "Data Summary"
        
        # Add table
        rows, cols = df.shape
        table = slide.shapes.add_table(rows + 1, cols, Inches(1), Inches(2), Inches(8), Inches(4)).table
        
        # Header row
        for col_idx, col_name in enumerate(df.columns):
            table.cell(0, col_idx).text = str(col_name)
        
        # Data rows
        for row_idx in range(rows):
            for col_idx in range(cols):
                table.cell(row_idx + 1, col_idx).text = str(df.iloc[row_idx, col_idx])
        
        prs.save(file_path)
    
    async def _export_parquet(self, df: pd.DataFrame, file_path: str, config: ExportConfig):
        """Export DataFrame as Parquet"""
        quality_settings = config.quality_settings or {}
        
        df.to_parquet(
            file_path,
            compression=quality_settings.get('compression', 'snappy'),
            index=quality_settings.get('include_index', False)
        )
    
    async def _export_html(self, df: pd.DataFrame, file_path: str, config: ExportConfig):
        """Export DataFrame as HTML"""
        quality_settings = config.quality_settings or {}
        
        html_string = df.to_html(
            index=quality_settings.get('include_index', False),
            escape=quality_settings.get('escape', False),
            table_id=quality_settings.get('table_id', 'data-table'),
            classes=quality_settings.get('css_classes', ['table', 'table-striped'])
        )
        
        # Add custom CSS if provided
        css = quality_settings.get('custom_css', '')
        if css:
            html_string = f"<style>{css}</style>\n{html_string}"
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(html_string)
    
    async def _export_markdown(self, df: pd.DataFrame, file_path: str, config: ExportConfig):
        """Export DataFrame as Markdown"""
        markdown_string = df.to_markdown(index=False)
        
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(markdown_string)
    
    async def _compress_file(self, file_path: str, compression: CompressionType) -> str:
        """Compress file using specified compression type"""
        if compression == CompressionType.GZIP:
            compressed_path = f"{file_path}.gz"
            with open(file_path, 'rb') as f_in:
                with gzip.open(compressed_path, 'wb') as f_out:
                    f_out.writelines(f_in)
            return compressed_path
        
        elif compression == CompressionType.ZIP:
            compressed_path = f"{file_path}.zip"
            with zipfile.ZipFile(compressed_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                zipf.write(file_path, file_path.split('/')[-1])
            return compressed_path
        
        elif compression == CompressionType.TAR:
            compressed_path = f"{file_path}.tar.gz"
            with tarfile.open(compressed_path, 'w:gz') as tar:
                tar.add(file_path, arcname=file_path.split('/')[-1])
            return compressed_path
        
        return file_path
    
    async def _encrypt_file(self, file_path: str) -> str:
        """Encrypt file using Fernet encryption"""
        if not self.cipher:
            raise ValueError("No encryption key provided")
        
        with open(file_path, 'rb') as f:
            file_data = f.read()
        
        encrypted_data = self.cipher.encrypt(file_data)
        encrypted_path = f"{file_path}.encrypted"
        
        with open(encrypted_path, 'wb') as f:
            f.write(encrypted_data)
        
        return encrypted_path
    
    async def _upload_to_destination(self, file_path: str, config: ExportConfig) -> str:
        """Upload file to specified destination"""
        destination_config = config.destination_config or {}
        
        if config.destination == StorageDestination.S3:
            return await self._upload_to_s3(file_path, destination_config)
        elif config.destination == StorageDestination.GCS:
            return await self._upload_to_gcs(file_path, destination_config)
        elif config.destination == StorageDestination.FTP:
            return await self._upload_to_ftp(file_path, destination_config)
        elif config.destination == StorageDestination.SFTP:
            return await self._upload_to_sftp(file_path, destination_config)
        
        return file_path
    
    async def _upload_to_s3(self, file_path: str, config: Dict[str, Any]) -> str:
        """Upload file to S3"""
        if not self.s3_client:
            self.s3_client = boto3.client(
                's3',
                aws_access_key_id=config.get('access_key'),
                aws_secret_access_key=config.get('secret_key'),
                region_name=config.get('region', 'us-east-1')
            )
        
        bucket = config['bucket']
        key = config.get('key', file_path.split('/')[-1])
        
        self.s3_client.upload_file(file_path, bucket, key)
        
        return f"s3://{bucket}/{key}"
    
    async def _upload_to_gcs(self, file_path: str, config: Dict[str, Any]) -> str:
        """Upload file to Google Cloud Storage"""
        if not self.gcs_client:
            self.gcs_client = gcs.Client(project=config.get('project_id'))
        
        bucket = self.gcs_client.bucket(config['bucket'])
        blob_name = config.get('blob_name', file_path.split('/')[-1])
        blob = bucket.blob(blob_name)
        
        blob.upload_from_filename(file_path)
        
        return f"gs://{config['bucket']}/{blob_name}"
    
    async def _upload_to_ftp(self, file_path: str, config: Dict[str, Any]) -> str:
        """Upload file to FTP server"""
        ftp = ftplib.FTP()
        ftp.connect(config['host'], config.get('port', 21))
        ftp.login(config['username'], config['password'])
        
        remote_path = config.get('remote_path', file_path.split('/')[-1])
        
        with open(file_path, 'rb') as f:
            ftp.storbinary(f'STOR {remote_path}', f)
        
        ftp.quit()
        
        return f"ftp://{config['host']}/{remote_path}"
    
    async def _upload_to_sftp(self, file_path: str, config: Dict[str, Any]) -> str:
        """Upload file to SFTP server"""
        ssh = paramiko.SSHClient()
        ssh.set_missing_host_key_policy(paramiko.AutoAddPolicy())
        ssh.connect(
            config['host'],
            port=config.get('port', 22),
            username=config['username'],
            password=config.get('password'),
            key_filename=config.get('key_file')
        )
        
        sftp = ssh.open_sftp()
        remote_path = config.get('remote_path', file_path.split('/')[-1])
        sftp.put(file_path, remote_path)
        
        sftp.close()
        ssh.close()
        
        return f"sftp://{config['host']}/{remote_path}"
    
    def _prepare_dataframe(self, data: Union[pd.DataFrame, List[Dict], Dict]) -> pd.DataFrame:
        """Convert various data types to DataFrame"""
        if isinstance(data, pd.DataFrame):
            return data
        elif isinstance(data, list):
            return pd.DataFrame(data)
        elif isinstance(data, dict):
            # Handle nested dictionaries
            if all(isinstance(v, (list, dict)) for v in data.values()):
                return pd.DataFrame(data)
            else:
                # Single record
                return pd.DataFrame([data])
        else:
            raise ValueError(f"Unsupported data type: {type(data)}")
    
    def _generate_file_path(self, config: ExportConfig) -> str:
        """Generate file path for export"""
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        
        if config.filename:
            filename = config.filename
        else:
            filename = f"export_{timestamp}"
        
        # Add extension based on format
        extensions = {
            ExportFormat.CSV: '.csv',
            ExportFormat.EXCEL: '.xlsx',
            ExportFormat.JSON: '.json',
            ExportFormat.XML: '.xml',
            ExportFormat.PDF: '.pdf',
            ExportFormat.POWERPOINT: '.pptx',
            ExportFormat.PARQUET: '.parquet',
            ExportFormat.HTML: '.html',
            ExportFormat.MARKDOWN: '.md',
            ExportFormat.TXT: '.txt'
        }
        
        extension = extensions.get(config.format, '')
        if not filename.endswith(extension):
            filename += extension
        
        return f"/tmp/{filename}"
    
    async def _get_file_size(self, file_path: str) -> int:
        """Get file size in bytes"""
        try:
            import os
            return os.path.getsize(file_path)
        except:
            return 0
    
    async def _calculate_checksum(self, file_path: str) -> str:
        """Calculate MD5 checksum of file"""
        import hashlib
        
        hash_md5 = hashlib.md5()
        try:
            with open(file_path, "rb") as f:
                for chunk in iter(lambda: f.read(4096), b""):
                    hash_md5.update(chunk)
            return hash_md5.hexdigest()
        except:
            return ""
    
    async def _add_excel_charts(self, worksheet, df: pd.DataFrame, quality_settings: Dict[str, Any]):
        """Add charts to Excel worksheet"""
        chart_configs = quality_settings.get('charts', [])
        
        for chart_config in chart_configs:
            chart_type = chart_config.get('type', 'column')
            data_range = chart_config.get('data_range', 'A1:B10')
            
            if chart_type == 'column':
                chart = BarChart()
            elif chart_type == 'line':
                chart = LineChart()
            elif chart_type == 'pie':
                chart = PieChart()
            else:
                continue
            
            data = Reference(worksheet, range_string=data_range)
            chart.add_data(data, titles_from_data=True)
            
            # Position chart
            chart_position = chart_config.get('position', 'D2')
            worksheet.add_chart(chart, chart_position)
    
    def generate_encryption_key(self) -> str:
        """Generate new encryption key"""
        return Fernet.generate_key().decode()
    
    async def decrypt_file(self, encrypted_file_path: str, output_path: str) -> bool:
        """Decrypt encrypted file"""
        if not self.cipher:
            raise ValueError("No encryption key provided")
        
        try:
            with open(encrypted_file_path, 'rb') as f:
                encrypted_data = f.read()
            
            decrypted_data = self.cipher.decrypt(encrypted_data)
            
            with open(output_path, 'wb') as f:
                f.write(decrypted_data)
            
            return True
        except Exception as e:
            print(f"Decryption failed: {str(e)}")
            return False